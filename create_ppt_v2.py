from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Load the original template
prs = Presentation(r'C:\Users\HP\Downloads\Code Catalyst 6.0_ppt.pptx')

# We have 6 slides from template, we need 7
# Add one more slide
blank_layout = prs.slide_layouts[1]  # Title and Content layout
prs.slides.add_slide(blank_layout)

# Colors
GREEN = (34, 120, 60)
DARK_GREEN = (20, 80, 40)
ORANGE = (255, 152, 0)
DARK_GRAY = (60, 60, 60)

def set_slide_content(slide, title_text, content_lines):
    """Set title and content for a slide"""
    # Set title
    title = slide.shapes.title
    if title:
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*GREEN)
    
    # Set content
    content = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            content = shape
            break
    
    if content:
        tf = content.text_frame
        tf.clear()
        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Check if line starts with special markers
            if line.startswith("##"):
                # Section header
                run = p.add_run()
                run.text = line.replace("##", "")
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*ORANGE)
            elif line.startswith(">>"):
                # Sub-point
                p.level = 1
                run = p.add_run()
                run.text = line.replace(">>", "  • ")
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(*DARK_GRAY)
            elif line.startswith("!!"):
                # Highlight
                run = p.add_run()
                run.text = line.replace("!!", "")
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*GREEN)
            else:
                run = p.add_run()
                run.text = line
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(*DARK_GRAY)
            
            p.space_after = Pt(4)

# ============================================
# SLIDE 1: TITLE PAGE (already exists - modify it)
# ============================================
slide1 = prs.slides[0]

# Clear existing and set new content
title1 = slide1.shapes.title
if title1:
    title1.text = "KisanSetu"
    for p in title1.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*GREEN)

subtitle1 = None
for shape in slide1.placeholders:
    if shape.placeholder_format.idx == 1:
        subtitle1 = shape
        break

if subtitle1:
    tf = subtitle1.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Farm to Consumer Digital Marketplace"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(*DARK_GRAY)
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "\nSmart India Hackathon 2026"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(*ORANGE)
    
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Problem Statement ID: SIH26033"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(*DARK_GRAY)
    
    p4 = tf.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "Ministry of Consumer Affairs, Food & Public Distribution"
    run4.font.size = Pt(12)
    run4.font.color.rgb = RGBColor(100, 100, 100)

# ============================================
# SLIDE 2: PROBLEM STATEMENT
# ============================================
slide2 = prs.slides[1]
set_slide_content(slide2, "Problem Statement (SIH26033)", [
    "##The Problem:",
    "Multiple intermediaries reduce farmers earnings and increase consumer prices.",
    "",
    "##Current Supply Chain Issues:",
    ">>Farmer gets only 30-40% of final price",
    ">>Consumer pays 40-50% more than fair price",
    ">>No direct communication between farmer & consumer",
    ">>No demand visibility for farmers",
    ">>Supply chain inefficiencies cause food waste",
    "",
    "##Expected Solution:",
    ">>Create digital marketplace connecting farmers/FPOs directly with consumers",
    ">>Provide logistics support with smart route optimization",
    ">>Use AI for demand forecasting and pricing"
])

# ============================================
# SLIDE 3: OUR SOLUTION
# ============================================
slide3 = prs.slides[2]
set_slide_content(slide3, "Our Solution: KisanSetu", [
    "##What is KisanSetu?",
    "A complete digital marketplace connecting farmers directly to consumers,",
    "eliminating middlemen and ensuring fair prices for both parties.",
    "",
    "##Key Features:",
    ">>Direct Marketplace: Farmers list, consumers buy - no middlemen",
    ">>AI Analytics: Demand forecasting, price suggestions, market trends",
    ">>Smart Logistics: Route optimization using Nearest Neighbor algorithm",
    ">>Real-time Chat: Direct communication via Socket.io",
    ">>Multi-role System: Farmer, Consumer, FPO, Bulk Buyer",
    "",
    "!!Impact: 30-40% more income for farmers, 20-30% lower prices for consumers"
])

# ============================================
# SLIDE 4: TECHNICAL APPROACH
# ============================================
slide4 = prs.slides[3]
set_slide_content(slide4, "Technical Approach", [
    "##Technology Stack (MERN):",
    ">>Frontend: React.js 18, Tailwind CSS, Framer Motion, Chart.js",
    ">>Backend: Node.js, Express.js, Socket.io",
    ">>Database: MongoDB Atlas (Cloud NoSQL)",
    ">>Auth: JWT tokens, bcrypt password hashing",
    ">>Deployment: Vercel (Frontend), Render (Backend)",
    "",
    "##System Architecture:",
    ">>Client Layer: React SPA with role-based interfaces",
    ">>Server Layer: Express.js REST API with middleware",
    ">>Database Layer: MongoDB with Mongoose ODM",
    "",
    "##Key Algorithms:",
    ">>Route Optimization: Nearest Neighbor TSP (O(n^2))",
    ">>Demand Forecasting: Time series analysis with moving averages"
])

# ============================================
# SLIDE 5: FEASIBILITY AND VIABILITY
# ============================================
slide5 = prs.slides[4]
set_slide_content(slide5, "Feasibility and Viability", [
    "##Feasibility:",
    ">>MERN stack is mature, well-documented, widely used",
    ">>MongoDB Atlas provides free cloud database tier",
    ">>Render & Vercel offer free hosting for deployment",
    ">>AI analytics built with proven algorithms",
    "",
    "##Challenges & Strategies:",
    ">>Challenge: Internet in rural areas → Strategy: Offline-first design",
    ">>Challenge: Building trust → Strategy: Rating system, verified profiles",
    ">>Challenge: Scaling → Strategy: Cloud-native architecture",
    "",
    "##Business Model:",
    ">>Commission-based: Small fee on successful transactions",
    ">>Premium subscriptions for advanced analytics",
    ">>Logistics partnerships for delivery services"
])

# ============================================
# SLIDE 6: IMPACT AND BENEFITS
# ============================================
slide6 = prs.slides[5]
set_slide_content(slide6, "Impact and Benefits", [
    "##Impact on Farmers:",
    ">>30-40% increase in earnings (direct sales)",
    ">>AI-powered demand forecasting for better planning",
    ">>Direct access to consumers and bulk buyers",
    "",
    "##Impact on Consumers:",
    ">>20-30% lower prices (no middlemen markup)",
    ">>Fresh produce directly from farms",
    ">>Transparent pricing and quality assurance",
    "",
    "##Impact on Society:",
    ">>Reduced food waste through demand forecasting",
    ">>Optimized delivery routes (less fuel consumption)",
    ">>Digital empowerment of small farmers",
    "",
    "!!Expected: 30% more farmer income, 25% less food waste, 35% faster delivery"
])

# ============================================
# SLIDE 7: FUTURE SCOPE & THANK YOU
# ============================================
slide7 = prs.slides[6]
set_slide_content(slide7, "Future Scope & Thank You", [
    "##Phase 2 Enhancements:",
    ">>Mobile app (React Native)",
    ">>Payment gateway integration (Razorpay)",
    ">>Multi-language support (Hindi, Tamil, etc.)",
    ">>GPS-based farmer discovery",
    "",
    "##Phase 3 Enhancements:",
    ">>IoT integration for crop monitoring",
    ">>Blockchain for supply chain transparency",
    ">>AI chatbot for farmer support",
    "",
    "##References:",
    ">>BigHaat, DeHaat, AgroStar (Existing platforms)",
    ">>FAO & NABARD reports on digital agriculture",
    "",
    "!!Thank You!",
    "KisanSetu - Bridging Farmers Directly to Consumers"
])

# Save the presentation
output_path = r'C:\Users\HP\Downloads\KisanSetu\KisanSetu_SIH_2026_Presentation.pptx'
prs.save(output_path)
print(f"PPT saved successfully at: {output_path}")
print(f"Total slides: {len(prs.slides)}")
