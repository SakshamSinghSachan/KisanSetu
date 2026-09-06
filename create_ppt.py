from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Load the template
prs = Presentation(r'C:\Users\HP\Downloads\Code Catalyst 6.0_ppt.pptx')

# Clear existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Helper functions
def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox

def add_shape_with_text(slide, left, top, width, height, text, font_size=14, bold=False, color=None, bg_color=None, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*bg_color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return shape

# Color scheme from template (green theme)
PRIMARY_GREEN = (34, 120, 60)
DARK_GREEN = (20, 80, 40)
LIGHT_GREEN = (220, 245, 220)
ORANGE = (255, 152, 0)
WHITE = (255, 255, 255)
BLACK = (0, 0, 0)
DARK_GRAY = (60, 60, 60)
LIGHT_GRAY = (240, 240, 240)

# ============================================
# SLIDE 1: TITLE PAGE
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[0])

# Title
title = slide1.shapes.title
title.text = "KisanSetu"
for paragraph in title.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

# Subtitle
subtitle = slide1.placeholders[1]
tf = subtitle.text_frame
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Farm to Consumer Digital Marketplace"
run.font.size = Pt(24)
run.font.color.rgb = RGBColor(*DARK_GRAY)

# Add more text
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "\nSmart India Hackathon 2026"
run2.font.size = Pt(20)
run2.font.color.rgb = RGBColor(*ORANGE)

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "Problem Statement ID: SIH26033"
run3.font.size = Pt(16)
run3.font.color.rgb = RGBColor(*DARK_GRAY)

p4 = tf.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run()
run4.text = "Ministry of Consumer Affairs, Food & Public Distribution"
run4.font.size = Pt(14)
run4.font.color.rgb = RGBColor(100, 100, 100)

# ============================================
# SLIDE 2: PROBLEM & SOLUTION
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[1])

# Title
title2 = slide2.shapes.title
title2.text = "Problem & Our Solution"
for paragraph in title2.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

# Content
content2 = slide2.placeholders[1]
tf2 = content2.text_frame
tf2.clear()

# Problem Section
p = tf2.paragraphs[0]
run = p.add_run()
run.text = "THE PROBLEM"
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = RGBColor(*ORANGE)

p2 = tf2.add_paragraph()
run2 = p2.add_run()
run2.text = "Multiple intermediaries reduce farmers earnings and increase consumer prices."
run2.font.size = Pt(14)
run2.font.color.rgb = RGBColor(*DARK_GRAY)

p3 = tf2.add_paragraph()
run3 = p3.add_run()
run3.text = "\nKey Issues:"
run3.font.size = Pt(16)
run3.font.bold = True
run3.font.color.rgb = RGBColor(*PRIMARY_GREEN)

problems = [
    "Farmer gets only 30-40% of final price",
    "Consumer pays 40-50% more than fair price",
    "No direct communication between farmer & consumer",
    "No demand visibility for farmers",
    "Supply chain inefficiencies cause food waste"
]
for prob in problems:
    pp = tf2.add_paragraph()
    pp.level = 1
    runp = pp.add_run()
    runp.text = f"• {prob}"
    runp.font.size = Pt(12)
    runp.font.color.rgb = RGBColor(*DARK_GRAY)

# Solution Section
p4 = tf2.add_paragraph()
p4.space_before = Pt(12)
run4 = p4.add_run()
run4.text = "OUR SOLUTION: KisanSetu"
run4.font.size = Pt(20)
run4.font.bold = True
run4.font.color.rgb = RGBColor(*PRIMARY_GREEN)

solutions = [
    "Direct farmer-to-consumer marketplace (no middlemen)",
    "AI-powered demand forecasting & price suggestions",
    "Smart logistics with route optimization",
    "Real-time chat between buyers and sellers",
    "Role-based access: Farmer, Consumer, FPO, Bulk Buyer"
]
for sol in solutions:
    ps = tf2.add_paragraph()
    runs = ps.add_run()
    runs.text = f"✓ {sol}"
    runs.font.size = Pt(12)
    runs.font.color.rgb = RGBColor(*DARK_GREEN)

# ============================================
# SLIDE 3: TECHNICAL APPROACH
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[1])

# Title
title3 = slide3.shapes.title
title3.text = "Technical Approach"
for paragraph in title3.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

# Content
content3 = slide3.placeholders[1]
tf3 = content3.text_frame
tf3.clear()

# Tech Stack
p = tf3.paragraphs[0]
run = p.add_run()
run.text = "TECHNOLOGY STACK (MERN)"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = RGBColor(*ORANGE)

tech_stack = [
    ("Frontend:", "React.js 18, Tailwind CSS, Framer Motion, Chart.js"),
    ("Backend:", "Node.js, Express.js, Socket.io"),
    ("Database:", "MongoDB Atlas (Cloud NoSQL)"),
    ("Authentication:", "JWT (JSON Web Tokens), bcrypt"),
    ("Deployment:", "Vercel (Frontend), Render (Backend)")
]
for label, value in tech_stack:
    pt = tf3.add_paragraph()
    runl = pt.add_run()
    runl.text = f"{label} "
    runl.font.size = Pt(12)
    runl.font.bold = True
    runl.font.color.rgb = RGBColor(*PRIMARY_GREEN)
    runv = pt.add_run()
    runv.text = value
    runv.font.size = Pt(12)
    runv.font.color.rgb = RGBColor(*DARK_GRAY)

# Key Features
p_feat = tf3.add_paragraph()
p_feat.space_before = Pt(12)
run_feat = p_feat.add_run()
run_feat.text = "KEY FEATURES"
run_feat.font.size = Pt(18)
run_feat.font.bold = True
run_feat.font.color.rgb = RGBColor(*ORANGE)

features = [
    "Multi-role Authentication (Farmer/Consumer/FPO/Bulk Buyer)",
    "Product Marketplace with search, filters, categories",
    "AI Demand Forecasting with price suggestions",
    "Smart Route Optimization (Nearest Neighbor Algorithm)",
    "Real-time Chat using Socket.io",
    "Analytics Dashboard with Chart.js visualization"
]
for feat in features:
    pf = tf3.add_paragraph()
    runf = pf.add_run()
    runf.text = f"→ {feat}"
    runf.font.size = Pt(12)
    runf.font.color.rgb = RGBColor(*DARK_GRAY)

# ============================================
# SLIDE 4: FLOWCHART / ARCHITECTURE
# ============================================
slide4 = prs.slides.add_slide(prs.slide_layouts[1])

# Title
title4 = slide4.shapes.title
title4.text = "System Architecture & Flow"
for paragraph in title4.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

# Content
content4 = slide4.placeholders[1]
tf4 = content4.text_frame
tf4.clear()

# Architecture
p = tf4.paragraphs[0]
run = p.add_run()
run.text = "SYSTEM ARCHITECTURE"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = RGBColor(*ORANGE)

arch_text = """
┌─────────────────────────────────────────────┐
│              CLIENT (React.js)              │
│   Farmer │ Consumer │ FPO │ Bulk Buyer      │
└────────────────────┬────────────────────────┘
                     │ REST API
┌────────────────────┴────────────────────────┐
│           SERVER (Node.js + Express)         │
│  Auth │ Products │ Orders │ Analytics │ Chat │
└────────────────────┬────────────────────────┘
                     │
┌────────────────────┴────────────────────────┐
│         DATABASE (MongoDB Atlas)             │
│  Users │ Products │ Orders │ Chats │ Messages│
└─────────────────────────────────────────────┘
"""
p_arch = tf4.add_paragraph()
run_arch = p_arch.add_run()
run_arch.text = arch_text
run_arch.font.size = Pt(10)
run_arch.font.name = "Courier New"
run_arch.font.color.rgb = RGBColor(*DARK_GRAY)

# Flow
p_flow = tf4.add_paragraph()
p_flow.space_before = Pt(8)
run_flow = p_flow.add_run()
run_flow.text = "USER FLOW"
run_flow.font.size = Pt(18)
run_flow.font.bold = True
run_flow.font.color.rgb = RGBColor(*ORANGE)

flow_steps = [
    "1. User registers (Farmer/Consumer) → JWT token generated",
    "2. Farmer lists products → Images, price, quantity",
    "3. Consumer browses → Search, filter, view details",
    "4. Consumer clicks Buy → Checkout with address",
    "5. Farmer receives order → Accepts/Rejects",
    "6. Smart logistics → Route optimized for delivery",
    "7. Order delivered → Payment & rating"
]
for step in flow_steps:
    ps = tf4.add_paragraph()
    runs = ps.add_run()
    runs.text = step
    runs.font.size = Pt(11)
    runs.font.color.rgb = RGBColor(*DARK_GRAY)

# ============================================
# SLIDE 5: FEASIBILITY AND VIABILITY
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[1])

# Title
title5 = slide5.shapes.title
title5.text = "Feasibility and Viability"
for paragraph in title5.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

# Content
content5 = slide5.placeholders[1]
tf5 = content5.text_frame
tf5.clear()

# Feasibility
p = tf5.paragraphs[0]
run = p.add_run()
run.text = "Feasibility"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

feasibility = [
    "MERN stack is mature, well-documented, and widely used",
    "MongoDB Atlas provides free cloud database tier",
    "Render & Vercel offer free hosting for deployment",
    "AI analytics can be built with existing algorithms",
    "Route optimization uses proven TSP algorithms"
]
for f in feasibility:
    pf = tf5.add_paragraph()
    runf = pf.add_run()
    runf.text = f"• {f}"
    runf.font.size = Pt(12)
    runf.font.color.rgb = RGBColor(*DARK_GRAY)

# Challenges
p2 = tf5.add_paragraph()
p2.space_before = Pt(10)
run2 = p2.add_run()
run2.text = "Challenges & Strategies"
run2.font.size = Pt(18)
run2.font.bold = True
run2.font.color.rgb = RGBColor(*ORANGE)

challenges = [
    ("Challenge:", "Internet connectivity in rural areas"),
    ("Strategy:", "Offline-first design with SMS fallback"),
    ("Challenge:", "Building trust between farmers and consumers"),
    ("Strategy:", "Rating system, verified profiles, secure payments"),
    ("Challenge:", "Scaling to millions of users"),
    ("Strategy:", "Cloud-native architecture, MongoDB Atlas auto-scaling")
]
for label, value in challenges:
    pc = tf5.add_paragraph()
    runl = pc.add_run()
    runl.text = f"{label} "
    runl.font.size = Pt(12)
    runl.font.bold = True
    runl.font.color.rgb = RGBColor(*DARK_GREEN) if label == "Strategy:" else RGBColor(*ORANGE)
    runv = pc.add_run()
    runv.text = value
    runv.font.size = Pt(12)
    runv.font.color.rgb = RGBColor(*DARK_GRAY)

# Business Model
p3 = tf5.add_paragraph()
p3.space_before = Pt(10)
run3 = p3.add_run()
run3.text = "Business Model"
run3.font.size = Pt(18)
run3.font.bold = True
run3.font.color.rgb = RGBColor(*PRIMARY_GREEN)

business = [
    "Commission-based: Small fee on successful transactions",
    "Premium subscriptions for advanced analytics",
    "Logistics partnerships for delivery services",
    "Government subsidies and agricultural schemes integration"
]
for b in business:
    pb = tf5.add_paragraph()
    runb = pb.add_run()
    runb.text = f"→ {b}"
    runb.font.size = Pt(12)
    runb.font.color.rgb = RGBColor(*DARK_GRAY)

# ============================================
# SLIDE 6: IMPACT AND BENEFITS
# ============================================
slide6 = prs.slides.add_slide(prs.slide_layouts[1])

# Title
title6 = slide6.shapes.title
title6.text = "Impact and Benefits"
for paragraph in title6.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

# Content
content6 = slide6.placeholders[1]
tf6 = content6.text_frame
tf6.clear()

# For Farmers
p = tf6.paragraphs[0]
run = p.add_run()
run.text = "Impact on Farmers"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

farmer_impact = [
    "30-40% increase in earnings (direct sales)",
    "AI-powered demand forecasting for better planning",
    "Direct access to consumers and bulk buyers",
    "Digital payment records and transparent pricing"
]
for f in farmer_impact:
    pf = tf6.add_paragraph()
    runf = pf.add_run()
    runf.text = f"✓ {f}"
    runf.font.size = Pt(12)
    runf.font.color.rgb = RGBColor(*DARK_GREEN)

# For Consumers
p2 = tf6.add_paragraph()
p2.space_before = Pt(10)
run2 = p2.add_run()
run2.text = "Impact on Consumers"
run2.font.size = Pt(18)
run2.font.bold = True
run2.font.color.rgb = RGBColor(*ORANGE)

consumer_impact = [
    "20-30% lower prices (no middlemen markup)",
    "Fresh produce directly from farms",
    "Transparent pricing and quality assurance",
    "Direct communication with farmers"
]
for c in consumer_impact:
    pc = tf6.add_paragraph()
    runc = pc.add_run()
    runc.text = f"✓ {c}"
    runc.font.size = Pt(12)
    runc.font.color.rgb = RGBColor(*DARK_GRAY)

# For Society
p3 = tf6.add_paragraph()
p3.space_before = Pt(10)
run3 = p3.add_run()
run3.text = "Impact on Society"
run3.font.size = Pt(18)
run3.font.bold = True
run3.font.color.rgb = RGBColor(*PRIMARY_GREEN)

society_impact = [
    "Reduced food waste through demand forecasting",
    "Optimized delivery routes (less fuel consumption)",
    "Digital empowerment of small farmers",
    "Transparent agricultural marketplace"
]
for s in society_impact:
    ps = tf6.add_paragraph()
    runs = ps.add_run()
    runs.text = f"✓ {s}"
    runs.font.size = Pt(12)
    runs.font.color.rgb = RGBColor(*DARK_GREEN)

# Key Metrics
p4 = tf6.add_paragraph()
p4.space_before = Pt(10)
run4 = p4.add_run()
run4.text = "Expected Metrics"
run4.font.size = Pt(18)
run4.font.bold = True
run4.font.color.rgb = RGBColor(*ORANGE)

metrics = [
    "Farmer Income: +30-40% increase",
    "Consumer Savings: 20-30% lower prices",
    "Food Waste Reduction: 25-30%",
    "Delivery Efficiency: 35% time saved"
]
for m in metrics:
    pm = tf6.add_paragraph()
    runm = pm.add_run()
    runm.text = f"📊 {m}"
    runm.font.size = Pt(12)
    runm.font.color.rgb = RGBColor(*DARK_GRAY)

# ============================================
# SLIDE 7: RESEARCH AND REFERENCES
# ============================================
slide7 = prs.slides.add_slide(prs.slide_layouts[1])

# Title
title7 = slide7.shapes.title
title7.text = "Research and References"
for paragraph in title7.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*PRIMARY_GREEN)

# Content
content7 = slide7.placeholders[1]
tf7 = content7.text_frame
tf7.clear()

# Existing Platforms
p = tf7.paragraphs[0]
run = p.add_run()
run.text = "Existing Platforms:"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = RGBColor(*ORANGE)

existing = [
    "BigHaat: https://bighaat.com (Agricultural e-commerce)",
    "DeHaat: https://dehaat.com (Farm-to-market platform)",
    "AgroStar: https://agrostar.in (Farm input marketplace)",
    "Ninjacart: https://ninjacart.com (B2B produce supply)"
]
for e in existing:
    pe = tf7.add_paragraph()
    rune = pe.add_run()
    rune.text = f"• {e}"
    rune.font.size = Pt(11)
    rune.font.color.rgb = RGBColor(*DARK_GRAY)

# Research Papers
p2 = tf7.add_paragraph()
p2.space_before = Pt(10)
run2 = p2.add_run()
run2.text = "Research Papers:"
run2.font.size = Pt(16)
run2.font.bold = True
run2.font.color.rgb = RGBColor(*PRIMARY_GREEN)

papers = [
    "FAO Report: Digital Agriculture for Smallholder Farmers",
    "NABARD: Impact of ICT on Agricultural Productivity",
    "IEEE: AI-Based Crop Recommendation Systems",
    "Springer: Route Optimization for Agricultural Supply Chain"
]
for paper in papers:
    pp = tf7.add_paragraph()
    runp = pp.add_run()
    runp.text = f"📄 {paper}"
    runp.font.size = Pt(11)
    runp.font.color.rgb = RGBColor(*DARK_GRAY)

# Comparison
p3 = tf7.add_paragraph()
p3.space_before = Pt(10)
run3 = p3.add_run()
run3.text = "Our Differentiation:"
run3.font.size = Pt(16)
run3.font.bold = True
run3.font.color.rgb = RGBColor(*PRIMARY_GREEN)

diff = [
    "Direct F2C model (no B2B intermediary layer)",
    "AI-powered demand forecasting built-in",
    "Smart logistics with route optimization",
    "Multi-role support (Farmer/Consumer/FPO/Buyer)",
    "Real-time chat for direct communication"
]
for d in diff:
    pd = tf7.add_paragraph()
    rund = pd.add_run()
    rund.text = f"★ {d}"
    rund.font.size = Pt(11)
    rund.font.color.rgb = RGBColor(*DARK_GREEN)

# Thank You
p4 = tf7.add_paragraph()
p4.space_before = Pt(16)
p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run()
run4.text = "Thank You!"
run4.font.size = Pt(28)
run4.font.bold = True
run4.font.color.rgb = RGBColor(*PRIMARY_GREEN)

p5 = tf7.add_paragraph()
p5.alignment = PP_ALIGN.CENTER
run5 = p5.add_run()
run5.text = "KisanSetu - Bridging Farmers Directly to Consumers"
run5.font.size = Pt(14)
run5.font.color.rgb = RGBColor(*DARK_GRAY)

# Save the presentation
output_path = r'C:\Users\HP\Downloads\KisanSetu\KisanSetu_SIH_Presentation.pptx'
prs.save(output_path)
print(f"PPT saved successfully at: {output_path}")
print(f"Total slides: {len(prs.slides)}")
