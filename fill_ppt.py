from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation(r'C:\Users\HP\Downloads\Code Catalyst 6.0_ppt.pptx')

def set_text(shape, lines, font_size=10, bold=False, color=None):
    shape.text_frame.clear()
    for i, line in enumerate(lines):
        p = shape.text_frame.paragraphs[0] if i == 0 else shape.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        if bold:
            run.font.bold = True
        if color:
            run.font.color.rgb = RGBColor(*color)

GREEN = (34, 120, 60)
ORANGE = (255, 152, 0)

# ===== SLIDE 1: TITLE =====
s1 = prs.slides[0]
for sh in s1.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text
    if "SMART INDIA HACKATHON" in t:
        set_text(sh, ["KISANSETU"], 40, True, GREEN)
    if "TITLE PAGE" in t:
        set_text(sh, ["Farm to Consumer Digital Marketplace"], 20)
    if "Problem Statement ID" in t:
        sh.text_frame.clear()
        lines = [
            "Problem Statement ID - SIH26033",
            "Problem Statement Title - Digital Marketplace for Farmers",
            "Theme - Agriculture, FoodTech & Rural Development",
            "PS Category - Software",
            "Team Name - KisanSetu Team"
        ]
        for i, line in enumerate(lines):
            p = sh.text_frame.paragraphs[0] if i == 0 else sh.text_frame.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(11)

# ===== SLIDE 2: PROBLEM & SOLUTION =====
s2 = prs.slides[1]
for sh in s2.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if "Smart Crop Advisory" in t:
        set_text(sh, ["Problem & Solution"], 28, True, GREEN)
    if "Proposed Solution" in t:
        set_text(sh, ["Our Solution: KisanSetu App"], 16, True, GREEN)
    if "Farmer Login" == t:
        set_text(sh, ["Farmer Registers"], 11, True)
    if "Choose the requirement" == t:
        set_text(sh, ["Lists Products"], 11, True)
    if "App analysis the requirement" == t:
        set_text(sh, ["Consumer Buys"], 11, True)
    if "Gives the advice by text, voice, or image" == t:
        set_text(sh, ["Direct Delivery"], 11, True)
    if "KrishiSakhi, a simple app" in t:
        lines = [
            "THE PROBLEM:",
            "Multiple intermediaries reduce farmers earnings",
            "and increase consumer prices.",
            "",
            "Key Issues:",
            "• Farmer gets only 30-40% of final price",
            "• Consumer pays 40-50% more",
            "• No direct farmer-consumer communication",
            "• No demand visibility for farmers",
            "• Supply chain causes food waste",
            "",
            "OUR SOLUTION: KisanSetu",
            "✓ Direct farmer-to-consumer marketplace",
            "✓ AI-powered demand forecasting",
            "✓ Smart logistics with route optimization",
            "✓ Real-time chat communication"
        ]
        set_text(sh, lines, 10)

# ===== SLIDE 3: TECHNICAL APPROACH =====
s3 = prs.slides[2]
for sh in s3.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t == "Frontend":
        set_text(sh, ["React.js + Tailwind"], 9, True, GREEN)
    if t == "APIs":
        set_text(sh, ["REST API"], 9, True, GREEN)
    if t == "AI / ML":
        set_text(sh, ["AI Analytics"], 9, True, GREEN)
    if t == "Backend":
        set_text(sh, ["Node.js + Express"], 9, True, GREEN)
    if t == "DevOps":
        set_text(sh, ["Vercel + Render"], 9, True, GREEN)
    if "Frontend: Flutter" in t:
        lines = [
            "TECHNOLOGY STACK (MERN)",
            "",
            "Frontend: React.js 18, Tailwind CSS, Framer Motion",
            "Backend: Node.js, Express.js, Socket.io",
            "Database: MongoDB Atlas (Cloud NoSQL)",
            "Auth: JWT tokens, bcrypt hashing",
            "Deployment: Vercel (Frontend), Render (Backend)",
            "",
            "KEY FEATURES:",
            "• Multi-role Auth (Farmer/Consumer/FPO)",
            "• Product Marketplace with search & filters",
            "• AI Demand Forecasting & price suggestions",
            "• Smart Route Optimization (Nearest Neighbor)",
            "• Real-time Chat via Socket.io",
            "• Analytics Dashboard with Chart.js"
        ]
        set_text(sh, lines, 10)

# ===== SLIDE 4: FEASIBILITY =====
s4 = prs.slides[3]
for sh in s4.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t.startswith("Feasibility") and "Cloud backend" in t:
        lines = [
            "Feasibility",
            "• MERN stack mature & widely used",
            "• MongoDB Atlas free tier",
            "• Render & Vercel free hosting",
            "• AI built with proven algorithms"
        ]
        set_text(sh, lines, 9, True, GREEN)
    if t.startswith("Challenges") and "Poor internet" in t:
        lines = [
            "Challenges",
            "• Internet in rural areas",
            "• Building user trust",
            "• Scaling to millions",
            "• Low digital literacy"
        ]
        set_text(sh, lines, 9, True, ORANGE)
    if t.startswith("Strategies") and "Offline-first" in t:
        lines = [
            "Strategies",
            "• Offline-first design",
            "• Rating system for trust",
            "• Cloud-native scaling",
            "• Multi-language support"
        ]
        set_text(sh, lines, 9, True, GREEN)
    if "Business Idea" in t and "Subscription-based" in t:
        lines = [
            "Business Model",
            "→ Commission on transactions",
            "→ Premium analytics subscription",
            "→ Logistics partnerships",
            "→ Government scheme integration"
        ]
        set_text(sh, lines, 9, True, ORANGE)

# ===== SLIDE 5: IMPACT =====
s5 = prs.slides[4]
for sh in s5.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t.startswith("Benefits") and "Yield Improvement" in t:
        lines = [
            "Benefits",
            "✓ Farmer Income: +30-40% increase",
            "✓ Consumer Savings: 20-30% lower prices",
            "✓ Food Waste: 25-30% reduction",
            "✓ Delivery: 35% time saved",
            "✓ Direct connection, no middlemen"
        ]
        set_text(sh, lines, 9, True, GREEN)
    if t.startswith("Impact") and "Economic Upliftment" in t:
        lines = [
            "Impact",
            "✓ Economic: Higher farmer income",
            "✓ Social: Digital empowerment",
            "✓ Environmental: Less food waste",
            "✓ Market: Transparent pricing",
            "✓ Scalable across India"
        ]
        set_text(sh, lines, 9, True, ORANGE)

# ===== SLIDE 6: REFERENCES =====
s6 = prs.slides[5]
for sh in s6.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t == "Existing Website & App:":
        set_text(sh, ["Existing Platforms:"], 10, True, GREEN)
    if t.startswith("AgroStar:") and "agrostar.in" in t:
        lines = [
            "Existing Platforms:",
            "• BigHaat: Agricultural e-commerce",
            "• DeHaat: Farm-to-market platform",
            "• AgroStar: Farm input marketplace",
            "• Ninjacart: B2B produce supply"
        ]
        set_text(sh, lines, 9, True)
    if t == "Research Innovation :":
        set_text(sh, ["Our Differentiation:"], 10, True, GREEN)
    if t.startswith("Farmer.chat:") and "arxiv" in t:
        lines = [
            "Our Differentiation:",
            "★ Direct F2C model (no B2B layer)",
            "★ AI demand forecasting built-in",
            "★ Smart logistics & route optimization",
            "★ Multi-role support system",
            "★ Real-time chat communication"
        ]
        set_text(sh, lines, 9, True)
    if t == "Comparision with Existing Systems:":
        set_text(sh, ["References:"], 10, True, GREEN)
    if t == "Refrences:":
        set_text(sh, ["Thank You! | KisanSetu"], 14, True, GREEN)

# Save
output = r'C:\Users\HP\Downloads\KisanSetu\KisanSetu_SIH_Presentation_Final.pptx'
prs.save(output)
print(f"PPT saved: {output}")
print(f"Total slides: {len(prs.slides)}")
